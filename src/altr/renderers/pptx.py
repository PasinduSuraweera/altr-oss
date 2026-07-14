from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .. import theme
from ..schemas import PresentationSpec, Slide

# Layout indexes in the standard template. Custom templates must keep the
# stock layout order (0 title, 1 title+content, 2 section header, 5 title only).
_TITLE_LAYOUT = 0
_BULLETS_LAYOUT = 1
_SECTION_LAYOUT = 2
_TITLE_ONLY_LAYOUT = 5

_CHART_TYPES = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
}

# Body area below the title on a 10 x 7.5 inch slide.
_BODY = (Inches(0.8), Inches(1.8), Inches(8.4), Inches(5.0))

# Bullet font size by indent level.
_BULLET_SIZES = (18, 16, 15, 14, 13)


def _rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def render_presentation(
    spec: PresentationSpec, out_dir: Path, template: Path | None = None
) -> Path:
    from . import output_path

    prs = Presentation(str(template)) if template else Presentation()
    styled = template is None  # never restyle a user-provided brand template

    for index, slide_spec in enumerate(spec.slides):
        add_slide(prs, slide_spec, is_first=index == 0, styled=styled)

    path = output_path(out_dir, spec.filename, ".pptx")
    prs.save(str(path))
    return path


def add_slide(prs, slide_spec: Slide, is_first: bool, styled: bool) -> None:
    """Append one slide to an open presentation (used by render and edit)."""
    layout = _effective_layout(slide_spec, is_first=is_first)
    slide = prs.slides.add_slide(prs.slide_layouts[_layout_index(layout)])
    slide.shapes.title.text = slide_spec.title

    if layout in ("title", "section") and slide_spec.subtitle:
        slide.placeholders[1].text = slide_spec.subtitle

    if layout == "bullets":
        body = slide.placeholders[1].text_frame
        for i, bullet in enumerate(slide_spec.bullets):
            para = body.paragraphs[0] if i == 0 else body.add_paragraph()
            para.text = bullet.text
            para.level = bullet.level

    if layout == "chart":
        _add_chart(slide, slide_spec, styled)

    if layout == "image":
        _add_image(slide, slide_spec)

    if slide_spec.notes:
        slide.notes_slide.notes_text_frame.text = slide_spec.notes

    if styled:
        _style_slide(slide, layout, slide_spec)


def _effective_layout(slide_spec: Slide, is_first: bool) -> str:
    """Repair layout choices local models commonly flub.

    A slide carrying chart data is a chart slide even if the model left the
    default 'bullets' layout; a body-less 'bullets' slide is really a title
    slide (deck opener) or a section divider.
    """
    layout = slide_spec.layout
    if layout == "bullets":
        if slide_spec.chart is not None:
            return "chart"
        if slide_spec.image is not None:
            return "image"
        if not slide_spec.bullets:
            return "title" if is_first else "section"
    return layout


def _layout_index(layout: str) -> int:
    if layout == "title":
        return _TITLE_LAYOUT
    if layout == "section":
        return _SECTION_LAYOUT
    if layout in ("chart", "image"):
        return _TITLE_ONLY_LAYOUT
    return _BULLETS_LAYOUT


# --- default theme styling ----------------------------------------------------


def _style_slide(slide, layout: str, slide_spec: Slide) -> None:
    if layout in ("title", "section"):
        _style_backdrop_slide(slide, layout, slide_spec)
    else:
        _style_content_slide(slide, layout, slide_spec)


def _style_backdrop_slide(slide, layout: str, slide_spec: Slide) -> None:
    """Title and section slides: white text on a deep accent backdrop."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(
        theme.ACCENT_DARK if layout == "title" else theme.ACCENT_DEEP
    )
    title_font = slide.shapes.title.text_frame.paragraphs[0].font
    title_font.size = Pt(40 if layout == "title" else 34)
    title_font.bold = True
    title_font.color.rgb = _rgb("FFFFFF")
    if slide_spec.subtitle:
        subtitle_font = slide.placeholders[1].text_frame.paragraphs[0].font
        subtitle_font.size = Pt(18)
        subtitle_font.color.rgb = _rgb(theme.ACCENT_TINT)


def _style_content_slide(slide, layout: str, slide_spec: Slide) -> None:
    """Bullets/chart/image slides: accent title over an accent rule, ink body."""
    title_para = slide.shapes.title.text_frame.paragraphs[0]
    title_para.alignment = PP_ALIGN.LEFT
    title_font = title_para.font
    title_font.size = Pt(28)
    title_font.bold = True
    title_font.color.rgb = _rgb(theme.ACCENT_DEEP)

    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.87), Inches(1.5), Inches(1.4), Pt(4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(theme.ACCENT)
    bar.line.fill.background()
    bar.shadow.inherit = False

    if layout == "bullets" and slide_spec.bullets:
        for para in slide.placeholders[1].text_frame.paragraphs:
            para.font.size = Pt(_BULLET_SIZES[min(para.level, len(_BULLET_SIZES) - 1)])
            para.font.color.rgb = _rgb(theme.INK)
            para.space_after = Pt(10)


def _style_chart(chart, kind: str, series_count: int) -> None:
    chart.has_title = False
    # A single bar/line series is named by the slide title; pie slices are
    # categories, so a pie always needs its legend for identity.
    chart.has_legend = series_count > 1 or kind == "pie"
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(12)
        chart.legend.font.color.rgb = _rgb(theme.INK_SECONDARY)

    if kind == "pie":
        try:  # per-slice colors; older python-pptx may not expose points
            for i, point in enumerate(chart.plots[0].series[0].points):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = _rgb(
                    theme.SERIES[i % len(theme.SERIES)]
                )
        except (AttributeError, NotImplementedError):
            pass
        return

    for i, series in enumerate(chart.series):
        color = _rgb(theme.SERIES[i % len(theme.SERIES)])
        if kind == "line":
            series.format.line.color.rgb = color
            series.format.line.width = Pt(2.25)
            series.smooth = False
        else:
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = color
            series.format.line.fill.background()

    for axis in (chart.category_axis, chart.value_axis):
        axis.tick_labels.font.size = Pt(12)
        axis.tick_labels.font.color.rgb = _rgb(theme.INK_MUTED)
        axis.format.line.color.rgb = _rgb(theme.GRIDLINE)
    chart.category_axis.has_major_gridlines = False
    chart.value_axis.has_major_gridlines = True
    gridline_format = chart.value_axis.major_gridlines.format.line
    gridline_format.color.rgb = _rgb(theme.GRIDLINE)
    gridline_format.width = Pt(0.75)


def _add_chart(slide, slide_spec: Slide, styled: bool) -> None:
    chart_spec = slide_spec.chart
    data = CategoryChartData()
    data.categories = chart_spec.categories
    for series in chart_spec.series:
        # Pad or trim so every series aligns with the category axis.
        values = list(series.values[: len(chart_spec.categories)])
        values += [0.0] * (len(chart_spec.categories) - len(values))
        data.add_series(series.name, values)
    frame = slide.shapes.add_chart(_CHART_TYPES[chart_spec.kind], *_BODY, data)
    if styled:
        _style_chart(frame.chart, chart_spec.kind, len(chart_spec.series))


def _add_image(slide, slide_spec: Slide) -> None:
    image_path = Path(slide_spec.image.path)
    if not image_path.is_file():
        raise FileNotFoundError(f"image not found: {slide_spec.image.path}")
    left, top, width, _ = _BODY
    slide.shapes.add_picture(str(image_path), left, top, width=width)
