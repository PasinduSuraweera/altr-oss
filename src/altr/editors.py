"""Read and edit existing Office files.

The editing tools mirror the create tools: the model sends a validated spec,
altr applies it. `read_office` gives the model the current content (with
paragraph/slide indexes) so its edits can reference real positions.

Editing is restricted to files inside the working directory or the output
directory, so the model can only touch documents the user is working with.
"""

from __future__ import annotations

import zipfile
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from .renderers.docx import _add_block
from .renderers.pptx import add_slide
from .renderers.xlsx import write_sheet
from .schemas import (
    AddSheet,
    AppendBlocks,
    AppendRows,
    AppendSlides,
    DeleteParagraph,
    EditDocumentSpec,
    EditPresentationSpec,
    EditSpreadsheetSpec,
    ReplaceSlideText,
    ReplaceText,
    SetCells,
    SetParagraph,
    SetSlideTitle,
)

_READ_ROW_CAP = 60  # rows per sheet shown to the model
_READ_TEXT_CAP = 300  # characters per paragraph/cell shown to the model


def resolve_existing(path_str: str, out_dir: Path) -> Path:
    """Resolve a model-supplied path, allowing only existing office files
    inside the working directory or the output directory."""
    path = Path(path_str).resolve()
    if path.suffix.lower() not in (".docx", ".xlsx", ".pptx"):
        raise ValueError(f"not an editable office file: {path_str}")
    if not path.is_file():
        raise FileNotFoundError(f"file not found: {path_str}")
    allowed = (Path.cwd().resolve(), Path(out_dir).resolve())
    if not any(path == base or base in path.parents for base in allowed):
        raise ValueError(
            f"refusing to touch {path_str}: outside the working and output directories"
        )
    return path


# --- read ---------------------------------------------------------------------


def read_office(path: Path) -> dict:
    """A compact JSON view of a file's content, sized for model context."""
    kind = path.suffix.lower()
    if kind == ".docx":
        return _read_docx(path)
    if kind == ".xlsx":
        return _read_xlsx(path)
    return _read_pptx(path)


def _clip(text: str) -> str:
    return text if len(text) <= _READ_TEXT_CAP else text[: _READ_TEXT_CAP] + "..."


def _read_docx(path: Path) -> dict:
    doc = Document(str(path))
    return {
        "type": "document",
        "paragraphs": [
            {"index": i, "style": p.style.name, "text": _clip(p.text)}
            for i, p in enumerate(doc.paragraphs)
        ],
        "tables": [
            {
                "index": i,
                "headers": [c.text for c in t.rows[0].cells],
                "data_rows": len(t.rows) - 1,
            }
            for i, t in enumerate(doc.tables)
        ],
    }


def _read_xlsx(path: Path) -> dict:
    wb = load_workbook(str(path))
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(max_row=_READ_ROW_CAP, values_only=True):
            rows.append([_clip(v) if isinstance(v, str) else v for v in row])
        sheets.append(
            {
                "name": ws.title,
                "rows_shown": len(rows),
                "total_rows": ws.max_row,
                "rows": rows,
            }
        )
    return {"type": "spreadsheet", "sheets": sheets, "has_charts": _has_charts(path)}


def _read_pptx(path: Path) -> dict:
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = [
            _clip(sh.text)
            for sh in slide.shapes
            if sh.has_text_frame and sh.text.strip()
        ]
        slides.append(
            {
                "index": i,
                "title": texts[0] if texts else "",
                "texts": texts[1:],
                "has_chart": any(sh.has_chart for sh in slide.shapes),
            }
        )
    return {"type": "presentation", "slides": slides}


def _has_charts(path: Path) -> bool:
    with zipfile.ZipFile(path) as z:
        return any("/charts/chart" in name for name in z.namelist())


# --- edit ---------------------------------------------------------------------


def edit_document(spec: EditDocumentSpec, path: Path) -> dict:
    doc = Document(str(path))
    for edit in spec.edits:
        if isinstance(edit, AppendBlocks):
            for block in edit.blocks:
                # styled=False: respect the existing document's own look
                _add_block(doc, block, styled=False)
        elif isinstance(edit, ReplaceText):
            for para in _all_paragraphs(doc):
                _replace_in_paragraph(para, edit.find, edit.replace)
        elif isinstance(edit, SetParagraph):
            _paragraph_at(doc, edit.index).text = edit.text
        elif isinstance(edit, DeleteParagraph):
            p = _paragraph_at(doc, edit.index)._element
            p.getparent().remove(p)
    doc.save(str(path))
    return {"ok": True, "file": str(path)}


def _paragraph_at(doc, index: int):
    paragraphs = doc.paragraphs
    if index >= len(paragraphs):
        raise ValueError(
            f"paragraph index {index} out of range (document has {len(paragraphs)})"
        )
    return paragraphs[index]


def _all_paragraphs(doc):
    yield from doc.paragraphs
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                yield from cell.paragraphs


def _replace_in_paragraph(para, find: str, replace: str) -> None:
    if find not in para.text:
        return
    for run in para.runs:  # run-level first: keeps formatting
        if find in run.text:
            run.text = run.text.replace(find, replace)
    if find in para.text:  # the match spans runs: rebuild the paragraph text
        para.text = para.text.replace(find, replace)


def edit_spreadsheet(spec: EditSpreadsheetSpec, path: Path) -> dict:
    had_charts = _has_charts(path)
    wb = load_workbook(str(path))
    for edit in spec.edits:
        if isinstance(edit, (SetCells, AppendRows)):
            if edit.sheet not in wb.sheetnames:
                raise ValueError(
                    f"no sheet named {edit.sheet!r} (has {wb.sheetnames})"
                )
            ws = wb[edit.sheet]
            if isinstance(edit, SetCells):
                for ref, value in edit.cells.items():
                    ws[ref] = value
            else:
                for row in edit.rows:
                    ws.append(row)
        elif isinstance(edit, AddSheet):
            write_sheet(wb, edit.sheet)
    wb.save(str(path))
    result = {"ok": True, "file": str(path)}
    if had_charts:
        result["warning"] = (
            "existing charts were dropped: the xlsx editor cannot preserve "
            "charts. Recreate the chart via add_sheet or create_spreadsheet, "
            "and tell the user."
        )
    return result


def edit_presentation(spec: EditPresentationSpec, path: Path) -> dict:
    prs = Presentation(str(path))
    for edit in spec.edits:
        if isinstance(edit, AppendSlides):
            for slide_spec in edit.slides:
                add_slide(prs, slide_spec, is_first=False, styled=True)
        elif isinstance(edit, SetSlideTitle):
            _slide_at(prs, edit.index).shapes.title.text = edit.title
        elif isinstance(edit, ReplaceSlideText):
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for para in shape.text_frame.paragraphs:
                        _replace_in_pptx_paragraph(para, edit.find, edit.replace)
    prs.save(str(path))
    return {"ok": True, "file": str(path)}


def _slide_at(prs, index: int):
    slides = list(prs.slides)
    if index >= len(slides):
        raise ValueError(f"slide index {index} out of range (deck has {len(slides)})")
    return slides[index]


def _replace_in_pptx_paragraph(para, find: str, replace: str) -> None:
    text = "".join(run.text for run in para.runs)
    if find not in text:
        return
    for run in para.runs:
        if find in run.text:
            run.text = run.text.replace(find, replace)
    text = "".join(run.text for run in para.runs)
    if find in text and para.runs:  # spans runs: collapse into the first run
        for run in para.runs[1:]:
            run.text = ""
        para.runs[0].text = text.replace(find, replace)
