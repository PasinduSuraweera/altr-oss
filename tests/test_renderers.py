from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from altr.renderers import output_path, render_document, render_presentation, render_spreadsheet
from altr.schemas import DocumentSpec, PresentationSpec, SpreadsheetSpec


def test_render_document_roundtrip(tmp_path):
    spec = DocumentSpec.model_validate(
        {
            "filename": "report.docx",
            "title": "Q3 Report",
            "blocks": [
                {"type": "heading", "level": 1, "text": "Summary"},
                {"type": "paragraph", "text": "Revenue grew 12%."},
                {"type": "bullet_list", "items": ["APAC up", "EMEA flat"]},
                {"type": "numbered_list", "items": ["Hire", "Ship"]},
                {"type": "table", "headers": ["Region", "Growth"], "rows": [["APAC", "20%"]]},
                {"type": "page_break"},
            ],
        }
    )
    path = render_document(spec, tmp_path)
    assert path == tmp_path / "report.docx"

    doc = Document(str(path))
    texts = [p.text for p in doc.paragraphs]
    assert "Q3 Report" in texts
    assert "Revenue grew 12%." in texts
    assert "APAC up" in texts
    table = doc.tables[0]
    assert table.rows[0].cells[0].text == "Region"
    assert table.rows[1].cells[1].text == "20%"


def test_render_spreadsheet_roundtrip(tmp_path):
    spec = SpreadsheetSpec.model_validate(
        {
            "filename": "budget",  # extension added automatically
            "sheets": [
                {
                    "name": "2026",
                    "columns": [{"header": "Item", "width": 24}, {"header": "Cost"}],
                    "rows": [["Servers", 1200], ["Domains", 40], ["Total", "=SUM(B2:B3)"]],
                }
            ],
        }
    )
    path = render_spreadsheet(spec, tmp_path)
    assert path.name == "budget.xlsx"

    wb = load_workbook(str(path))
    ws = wb["2026"]
    assert ws["A1"].value == "Item"
    assert ws["A1"].font.bold
    assert ws["B2"].value == 1200
    assert ws["B4"].value == "=SUM(B2:B3)"
    assert ws.freeze_panes == "A2"


def test_render_presentation_roundtrip(tmp_path):
    spec = PresentationSpec.model_validate(
        {
            "filename": "pitch.pptx",
            "slides": [
                {"layout": "title", "title": "SolarDrones", "subtitle": "Series A"},
                {
                    "layout": "bullets",
                    "title": "Why now",
                    "bullets": [{"text": "Costs down 40%"}, {"text": "Detail", "level": 1}],
                    "notes": "Pause here.",
                },
                {"layout": "section", "title": "Financials"},
            ],
        }
    )
    path = render_presentation(spec, tmp_path)

    prs = Presentation(str(path))
    slides = list(prs.slides)
    assert len(slides) == 3
    assert slides[0].shapes.title.text == "SolarDrones"
    assert slides[1].shapes.title.text == "Why now"
    assert "Costs down 40%" in slides[1].placeholders[1].text_frame.text
    assert slides[1].notes_slide.notes_text_frame.text == "Pause here."


def test_output_path_blocks_traversal(tmp_path):
    path = output_path(tmp_path, "../../etc/passwd", ".docx")
    assert path.parent == tmp_path
    assert path.name == "passwd.docx"


def test_string_bullets_are_accepted(tmp_path):
    spec = PresentationSpec.model_validate(
        {
            "filename": "pitch.pptx",
            "slides": [
                {
                    "layout": "bullets",
                    "title": "Why altr",
                    "bullets": ["Plain string", {"text": "Object bullet", "level": 1}],
                }
            ],
        }
    )
    assert [(b.text, b.level) for b in spec.slides[0].bullets] == [
        ("Plain string", 0),
        ("Object bullet", 1),
    ]
    path = render_presentation(spec, tmp_path)
    prs = Presentation(str(path))
    texts = [
        p.text
        for shape in prs.slides[0].shapes
        if shape.has_text_frame
        for p in shape.text_frame.paragraphs
    ]
    assert "Plain string" in texts
    assert "Object bullet" in texts


def test_chart_without_value_columns_plots_numeric_columns(tmp_path):
    spec = SpreadsheetSpec.model_validate(
        {
            "filename": "costs.xlsx",
            "sheets": [
                {
                    "name": "Costs",
                    "columns": [{"header": "Item"}, {"header": "Cost"}],
                    "rows": [["Servers", 1200], ["Domains", 40]],
                    "charts": [{"kind": "bar", "title": "Costs"}],
                }
            ],
        }
    )
    path = render_spreadsheet(spec, tmp_path)
    wb = load_workbook(str(path))
    charts = wb["Costs"]._charts
    assert len(charts) == 1
    # The lone numeric column (B) was picked up as the data series.
    assert "B" in charts[0].series[0].val.numRef.f


def test_table_cells_render_inline_markdown(tmp_path):
    spec = DocumentSpec.model_validate(
        {
            "filename": "cells.docx",
            "blocks": [
                {
                    "type": "table",
                    "headers": ["**Benchmark**", "Notes"],
                    "rows": [["**MMLU**", "*multi-subject* QA"]],
                }
            ],
        }
    )
    path = render_document(spec, tmp_path)
    table = Document(str(path)).tables[0]
    assert table.rows[0].cells[0].text == "Benchmark"  # no literal asterisks
    body_cell = table.rows[1].cells[0]
    assert body_cell.text == "MMLU"
    assert body_cell.paragraphs[0].runs[0].bold
    notes_runs = table.rows[1].cells[1].paragraphs[0].runs
    assert any(r.italic for r in notes_runs)
    assert "*" not in table.rows[1].cells[1].text
