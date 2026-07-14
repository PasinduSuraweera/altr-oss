"""Editing tests: create a file with the renderers, edit it, re-open, assert."""

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from altr.renderers import render_document, render_presentation, render_spreadsheet
from altr.schemas import DocumentSpec, PresentationSpec, SpreadsheetSpec
from altr.tools import dispatch


def _make_doc(tmp_path):
    spec = DocumentSpec.model_validate(
        {
            "filename": "notes.docx",
            "title": "Notes",
            "blocks": [
                {"type": "paragraph", "text": "Project Foo kicks off in March."},
                {"type": "paragraph", "text": "Owner: TBD"},
            ],
        }
    )
    return render_document(spec, tmp_path)


def test_edit_document_ops(tmp_path):
    path = _make_doc(tmp_path)
    outcome = dispatch(
        "edit_document",
        {
            "path": str(path),
            "edits": [
                {"op": "replace_text", "find": "Foo", "replace": "Phoenix"},
                {"op": "set_paragraph", "index": 2, "text": "Owner: Priya"},
                {
                    "op": "append_blocks",
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Risks"},
                        {"type": "bullet_list", "items": ["Scope creep"]},
                    ],
                },
            ],
        },
        out_dir=tmp_path,
    )
    assert outcome["ok"], outcome
    texts = [p.text for p in Document(str(path)).paragraphs]
    assert "Project Phoenix kicks off in March." in texts
    assert "Owner: Priya" in texts
    assert "Risks" in texts
    assert "Scope creep" in texts


def test_delete_paragraph(tmp_path):
    path = _make_doc(tmp_path)
    outcome = dispatch(
        "edit_document",
        {"path": str(path), "edits": [{"op": "delete_paragraph", "index": 2}]},
        out_dir=tmp_path,
    )
    assert outcome["ok"]
    assert "Owner: TBD" not in [p.text for p in Document(str(path)).paragraphs]


def test_edit_spreadsheet_ops(tmp_path):
    spec = SpreadsheetSpec.model_validate(
        {
            "filename": "budget.xlsx",
            "sheets": [
                {
                    "name": "Costs",
                    "columns": [{"header": "Item"}, {"header": "Cost"}],
                    "rows": [["Hosting", 30]],
                }
            ],
        }
    )
    path = render_spreadsheet(spec, tmp_path)
    outcome = dispatch(
        "edit_spreadsheet",
        {
            "path": str(path),
            "edits": [
                {"op": "set_cells", "sheet": "Costs", "cells": {"B2": 45}},
                {"op": "append_rows", "sheet": "Costs", "rows": [["Domains", 12]]},
                {
                    "op": "add_sheet",
                    "sheet": {"name": "Notes", "rows": [["reviewed", True]]},
                },
            ],
        },
        out_dir=tmp_path,
    )
    assert outcome["ok"], outcome
    wb = load_workbook(str(path))
    assert wb["Costs"]["B2"].value == 45
    assert wb["Costs"]["A3"].value == "Domains"
    assert wb["Notes"]["A1"].value == "reviewed"


def test_edit_spreadsheet_warns_about_charts(tmp_path):
    spec = SpreadsheetSpec.model_validate(
        {
            "filename": "charted.xlsx",
            "sheets": [
                {
                    "name": "Data",
                    "columns": [{"header": "K"}, {"header": "V"}],
                    "rows": [["a", 1], ["b", 2]],
                    "charts": [{"kind": "bar", "title": "V"}],
                }
            ],
        }
    )
    path = render_spreadsheet(spec, tmp_path)
    outcome = dispatch(
        "edit_spreadsheet",
        {"path": str(path), "edits": [{"op": "set_cells", "sheet": "Data", "cells": {"B2": 9}}]},
        out_dir=tmp_path,
    )
    assert outcome["ok"]
    assert "charts" in outcome.get("warning", "")


def test_edit_presentation_ops(tmp_path):
    spec = PresentationSpec.model_validate(
        {
            "filename": "deck.pptx",
            "slides": [
                {"layout": "title", "title": "Q3 Review", "subtitle": "Draft"},
                {"layout": "bullets", "title": "Wins", "bullets": ["Shipped v1"]},
            ],
        }
    )
    path = render_presentation(spec, tmp_path)
    outcome = dispatch(
        "edit_presentation",
        {
            "path": str(path),
            "edits": [
                {"op": "set_slide_title", "index": 0, "title": "Q3 Business Review"},
                {"op": "replace_text", "find": "Shipped v1", "replace": "Shipped v2"},
                {
                    "op": "append_slides",
                    "slides": [{"layout": "bullets", "title": "Next", "bullets": ["Plan Q4"]}],
                },
            ],
        },
        out_dir=tmp_path,
    )
    assert outcome["ok"], outcome
    prs = Presentation(str(path))
    assert len(prs.slides) == 3
    all_text = "\n".join(
        sh.text for s in prs.slides for sh in s.shapes if sh.has_text_frame
    )
    assert "Q3 Business Review" in all_text
    assert "Shipped v2" in all_text
    assert "Plan Q4" in all_text


def test_read_office_file(tmp_path):
    path = _make_doc(tmp_path)
    outcome = dispatch("read_office_file", {"path": str(path)}, out_dir=tmp_path)
    assert outcome["ok"]
    content = outcome["content"]
    assert content["type"] == "document"
    assert any("Project Foo" in p["text"] for p in content["paragraphs"])
    assert all("index" in p for p in content["paragraphs"])


def test_edit_rejects_paths_outside_workspace(tmp_path):
    outcome = dispatch(
        "read_office_file",
        {"path": "C:/Windows/whatever.docx"},
        out_dir=tmp_path,
    )
    assert outcome["ok"] is False


def test_edit_rejects_missing_and_non_office_files(tmp_path):
    missing = dispatch(
        "read_office_file", {"path": str(tmp_path / "nope.docx")}, out_dir=tmp_path
    )
    assert missing["ok"] is False and "not found" in missing["error"]
    (tmp_path / "raw.txt").write_text("hi")
    wrong = dispatch(
        "read_office_file", {"path": str(tmp_path / "raw.txt")}, out_dir=tmp_path
    )
    assert wrong["ok"] is False


def test_mcp_tool_definitions():
    mcp = pytest.importorskip("mcp")  # noqa: F841 - optional dependency
    from altr.mcp_server import build_server
    from altr.tools import get_tool_specs

    build_server("output")  # must construct without error
    names = [name for name, _, _ in get_tool_specs()]
    assert "create_document" in names
    assert "edit_document" in names
    assert "read_office_file" in names
    assert len(names) == 7
