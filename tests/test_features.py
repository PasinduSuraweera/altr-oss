"""Tests for charts, images, markdown blocks, templates, and PDF export."""

import base64

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from altr.pdf import soffice_available, to_pdf
from altr.renderers import render_document, render_presentation, render_spreadsheet
from altr.schemas import DocumentSpec, PresentationSpec, SpreadsheetSpec

# A 1x1 transparent PNG.
_PNG = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8BQDwAEhQGAhKmMIQAAAABJRU5ErkJggg=="
)


@pytest.fixture
def png(tmp_path):
    path = tmp_path / "pixel.png"
    path.write_bytes(_PNG)
    return path


def test_spreadsheet_chart(tmp_path):
    spec = SpreadsheetSpec.model_validate(
        {
            "filename": "sales.xlsx",
            "sheets": [
                {
                    "name": "Q1",
                    "columns": [{"header": "Month"}, {"header": "Revenue"}],
                    "rows": [["Jan", 100], ["Feb", 140], ["Mar", 90]],
                    "charts": [
                        {
                            "kind": "bar",
                            "title": "Revenue by month",
                            "label_column": 1,
                            "value_columns": [2],
                        }
                    ],
                }
            ],
        }
    )
    path = render_spreadsheet(spec, tmp_path)
    ws = load_workbook(str(path))["Q1"]
    assert len(ws._charts) == 1
    assert ws._charts[0].title is not None


def test_slide_chart(tmp_path):
    spec = PresentationSpec.model_validate(
        {
            "filename": "metrics.pptx",
            "slides": [
                {
                    "layout": "chart",
                    "title": "Growth",
                    "chart": {
                        "kind": "line",
                        "categories": ["Q1", "Q2", "Q3"],
                        "series": [{"name": "Users", "values": [10, 25, 60]}],
                    },
                }
            ],
        }
    )
    path = render_presentation(spec, tmp_path)
    slide = list(Presentation(str(path)).slides)[0]
    assert any(shape.has_chart for shape in slide.shapes)


def test_chart_slide_requires_chart():
    with pytest.raises(ValueError):
        PresentationSpec.model_validate(
            {"filename": "x.pptx", "slides": [{"layout": "chart", "title": "no data"}]}
        )


def test_document_image(tmp_path, png):
    spec = DocumentSpec.model_validate(
        {
            "filename": "illustrated.docx",
            "blocks": [
                {"type": "image", "path": str(png), "width_inches": 2, "caption": "A pixel"},
            ],
        }
    )
    path = render_document(spec, tmp_path)
    doc = Document(str(path))
    assert len(doc.inline_shapes) == 1
    assert any(p.text == "A pixel" for p in doc.paragraphs)


def test_document_image_missing_file(tmp_path):
    spec = DocumentSpec.model_validate(
        {
            "filename": "broken.docx",
            "blocks": [{"type": "image", "path": str(tmp_path / "nope.png")}],
        }
    )
    with pytest.raises(FileNotFoundError):
        render_document(spec, tmp_path)


def test_slide_image(tmp_path, png):
    spec = PresentationSpec.model_validate(
        {
            "filename": "photo.pptx",
            "slides": [
                {"layout": "image", "title": "The pixel", "image": {"path": str(png)}}
            ],
        }
    )
    path = render_presentation(spec, tmp_path)
    slide = list(Presentation(str(path)).slides)[0]
    assert any(shape.shape_type == 13 for shape in slide.shapes)  # 13 = PICTURE


def test_markdown_block_renders(tmp_path):
    spec = DocumentSpec.model_validate(
        {
            "filename": "md.docx",
            "blocks": [
                {"type": "markdown", "text": "# Section\n\nSome **bold** text.\n\n- a\n- b"}
            ],
        }
    )
    path = render_document(spec, tmp_path)
    doc = Document(str(path))
    texts = [p.text for p in doc.paragraphs]
    assert "Section" in texts
    assert "Some bold text." in texts
    bold_para = next(p for p in doc.paragraphs if p.text == "Some bold text.")
    assert any(run.bold for run in bold_para.runs)


def test_docx_template(tmp_path):
    template = tmp_path / "brand.docx"
    tdoc = Document()
    tdoc.add_paragraph("CONFIDENTIAL - ACME CORP")
    tdoc.save(str(template))

    spec = DocumentSpec.model_validate(
        {"filename": "memo.docx", "blocks": [{"type": "paragraph", "text": "hello"}]}
    )
    path = render_document(spec, tmp_path / "out", template=template)
    texts = [p.text for p in Document(str(path)).paragraphs]
    assert "CONFIDENTIAL - ACME CORP" in texts
    assert "hello" in texts


def test_pptx_template(tmp_path):
    template = tmp_path / "brand.pptx"
    Presentation().save(str(template))

    spec = PresentationSpec.model_validate(
        {"filename": "deck.pptx", "slides": [{"layout": "bullets", "title": "Hi"}]}
    )
    path = render_presentation(spec, tmp_path / "out", template=template)
    assert list(Presentation(str(path)).slides)[0].shapes.title.text == "Hi"


@pytest.mark.skipif(not soffice_available(), reason="LibreOffice not installed")
def test_pdf_export(tmp_path):
    spec = DocumentSpec.model_validate(
        {"filename": "print-me.docx", "blocks": [{"type": "paragraph", "text": "pdf me"}]}
    )
    docx_path = render_document(spec, tmp_path)
    pdf_path = to_pdf(docx_path)
    assert pdf_path.is_file()
    assert pdf_path.read_bytes()[:5] == b"%PDF-"
